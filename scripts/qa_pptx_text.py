from __future__ import annotations

from pathlib import Path

from pptx import Presentation


PPTX_PATH = Path(__file__).resolve().parents[1] / "reports" / "predictive_maintenance_lg_summary.pptx"


def main() -> None:
    prs = Presentation(PPTX_PATH)
    print(f"slides={len(prs.slides)}")
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip().replace("\n", " | "))
        joined = " || ".join(texts)
        print(f"SLIDE {idx}: {joined[:1000]}")
        for token in ["xxxx", "lorem", "ipsum"]:
            if token in joined.lower():
                raise SystemExit(f"placeholder token found on slide {idx}: {token}")


if __name__ == "__main__":
    main()
