from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_ROOT = Path(__file__).resolve().parents[1]
REPORTS_DIR = PROJECT_ROOT / "reports"
OUTPUT_DIR = PROJECT_ROOT / "outputs"
PPTX_PATH = REPORTS_DIR / "predictive_maintenance_lg_summary.pptx"

SUMMARY_PATH = OUTPUT_DIR / "experiments" / "predictive_maintenance" / "experiment_summary.csv"
OOF_PATH = OUTPUT_DIR / "experiments" / "predictive_maintenance" / "oof_predictions.csv"
DATA_PATH = OUTPUT_DIR / "data_gathering" / "predictive_maintenance" / "processed_predictive_maintenance.csv"
FEATURE_PATH = OUTPUT_DIR / "data_gathering" / "predictive_maintenance" / "feature_importance.csv"
CONFUSION_PATH = OUTPUT_DIR / "data_gathering" / "predictive_maintenance" / "confusion_matrix.png"
PR_PATH = OUTPUT_DIR / "data_gathering" / "predictive_maintenance" / "precision_recall_curve.png"
FEATURE_IMG_PATH = OUTPUT_DIR / "data_gathering" / "predictive_maintenance" / "feature_importance.png"
TARGET_IMG_PATH = OUTPUT_DIR / "data_gathering" / "predictive_maintenance" / "target_distribution.png"

LG_RED = RGBColor(165, 0, 52)
LG_DARK = RGBColor(38, 38, 38)
LG_GRAY = RGBColor(102, 102, 102)
LG_LIGHT = RGBColor(247, 244, 244)
WHITE = RGBColor(255, 255, 255)
PINK = RGBColor(242, 211, 219)


def add_bg(slide, color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.z_order = 0


def add_lg_mark(slide, x=0.35, y=0.25, size=0.52, inverse=False):
    fill = WHITE if inverse else LG_RED
    text_color = LG_RED if inverse else WHITE
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    circle.fill.solid()
    circle.fill.fore_color.rgb = fill
    circle.line.color.rgb = fill
    tx = circle.text_frame
    tx.clear()
    p = tx.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "LG"
    run.font.bold = True
    run.font.size = Pt(13)
    run.font.color.rgb = text_color
    return circle


def add_title(slide, title, subtitle=None, dark=False):
    add_lg_mark(slide, inverse=dark)
    color = WHITE if dark else LG_DARK
    box = slide.shapes.add_textbox(Inches(1.0), Inches(0.28), Inches(11.6), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(25)
    run.font.bold = True
    run.font.color.rgb = color
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(1.02), Inches(0.83), Inches(11.2), Inches(0.35))
        stf = sbox.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(10.5)
        sr.font.color.rgb = PINK if dark else LG_GRAY


def add_footer(slide, page):
    box = slide.shapes.add_textbox(Inches(11.7), Inches(7.05), Inches(1.2), Inches(0.25))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{page:02d}"
    r.font.size = Pt(9)
    r.font.color.rgb = LG_GRAY


def add_text(slide, text, x, y, w, h, size=14, color=LG_DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, title, value, caption=None, fill=LG_LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(230, 220, 224)
    tf = shape.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(24)
    r1.font.bold = True
    r1.font.color.rgb = LG_RED
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = title
    r2.font.size = Pt(11)
    r2.font.bold = True
    r2.font.color.rgb = LG_DARK
    if caption:
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = caption
        r3.font.size = Pt(8.5)
        r3.font.color.rgb = LG_GRAY


def add_bullets(slide, items, x, y, w, h, size=13, color=LG_DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Arial"
        p.space_after = Pt(5)
    return box


def add_table(slide, df, x, y, w, h, font_size=8):
    rows, cols = df.shape[0] + 1, df.shape[1]
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = LG_RED
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE
                r.font.bold = True
                r.font.size = Pt(font_size)
    for i, (_, row) in enumerate(df.iterrows(), start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else LG_LIGHT
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = LG_DARK
                    r.font.size = Pt(font_size)
    return table_shape


def short_feature_label(name: str) -> str:
    mapping = {
        "power_proxy": "Power proxy",
        "Tool wear [min]": "Tool wear",
        "Rotational speed [rpm]": "RPM",
        "Torque [Nm]": "Torque",
        "temperature_gap": "Temp gap",
        "wear_per_torque": "Wear/Torque",
        "Process temperature [K]": "Process temp",
        "Air temperature [K]": "Air temp",
    }
    return mapping.get(name, name[:18])


def create_deck():
    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    df = pd.read_csv(DATA_PATH)
    summary = pd.read_csv(SUMMARY_PATH)
    feature = pd.read_csv(FEATURE_PATH)
    oof = pd.read_csv(OOF_PATH)
    best = summary.iloc[0]
    best_exp = best["experiment_id"]
    best_oof = oof[oof["experiment_id"] == best_exp].copy()
    errors = best_oof[best_oof["actual"] != best_oof["predicted_label"]].copy()
    fn = int(((best_oof["actual"] == 1) & (best_oof["predicted_label"] == 0)).sum())
    fp = int(((best_oof["actual"] == 0) & (best_oof["predicted_label"] == 1)).sum())
    near = errors.assign(margin=(errors["predicted_probability"] - 0.5).abs()).sort_values("margin").head(3)
    high = errors.assign(margin=(errors["predicted_probability"] - 0.5).abs()).sort_values("margin", ascending=False).head(3)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_bg(slide, LG_RED)
    add_lg_mark(slide, x=0.65, y=0.55, size=0.72, inverse=True)
    add_text(slide, "Machine Predictive\nMaintenance", 0.85, 1.55, 6.3, 1.45, size=36, color=WHITE, bold=True)
    add_text(slide, "설비 센서 기반 고장 예측 모델링 · 비교실험 · 오분류 분석", 0.9, 3.15, 7.0, 0.35, size=15, color=PINK)
    add_card(slide, 8.1, 1.1, 1.35, 1.25, "rows", "10K", "tabular")
    add_card(slide, 9.65, 1.1, 1.35, 1.25, "failure rate", "3.39%", "imbalanced")
    add_card(slide, 11.2, 1.1, 1.35, 1.25, "best model", "LGBM", "5-fold CV")
    add_text(slide, "LG-inspired theme: LG Red · Charcoal · White", 0.9, 6.75, 5.0, 0.25, size=10, color=PINK)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    add_title(slide, "1. 데이터셋 소개", "무엇을 활용해 무엇을 예측하는가")
    add_text(slide, "Kaggle: shivamb/machine-predictive-maintenance-classification", 0.9, 1.25, 6.0, 0.3, size=11, color=LG_GRAY)
    add_bullets(
        slide,
        [
            "센서/운전 조건: 온도, 회전속도, 토크, 공구 마모, 제품 타입",
            "목표: 기계 고장 여부(Target=1)를 사전에 예측",
            "고장 유형(Failure Type)은 2차 분석용이며, 1차 모델 feature에서는 제외",
            "고장 비율이 3.39%라 accuracy보다 recall, F1, PR-AUC가 중요",
        ],
        0.9,
        1.75,
        5.6,
        2.4,
        size=13,
    )
    slide.shapes.add_picture(str(TARGET_IMG_PATH), Inches(7.0), Inches(1.35), width=Inches(5.2), height=Inches(3.4))
    add_card(slide, 0.9, 4.75, 1.6, 1.05, "features", "8", "engineered")
    add_card(slide, 2.7, 4.75, 1.6, 1.05, "target", "binary", "failure")
    add_card(slide, 4.5, 4.75, 1.6, 1.05, "mismatch", "27", "label review")
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_bg(slide, LG_LIGHT)
    add_title(slide, "2. 실험 설계와 최고 성능 세팅", "baseline에서 CV 비교실험으로 점진 강화")
    top = summary.head(5)[["experiment_id", "model_name", "stage", "f1_mean", "pr_auc_mean", "recall_mean"]].copy()
    for col in ["f1_mean", "pr_auc_mean", "recall_mean"]:
        top[col] = top[col].map(lambda x: f"{x:.3f}")
    add_table(slide, top, 0.65, 1.25, 7.0, 2.45, font_size=8)
    add_card(slide, 8.2, 1.25, 1.55, 1.0, "Best exp", str(best_exp), "rank #1")
    add_card(slide, 9.95, 1.25, 1.55, 1.0, "F1 mean", f"{best['f1_mean']:.3f}", "5-fold")
    add_card(slide, 11.7, 1.25, 1.15, 1.0, "PR-AUC", f"{best['pr_auc_mean']:.3f}", "")
    add_bullets(
        slide,
        [
            "Model: LightGBM tuned deeper",
            "Feature: all engineered features included",
            "Imbalance: scale_pos_weight 적용",
            "Missing: median/mode 기준, Outlier: 별도 clipping 없음",
            "핵심 해석: 복잡도를 늘리되 과도한 sampling보다 안정적 boosting이 유리",
        ],
        8.2,
        2.65,
        4.5,
        2.6,
        size=12,
    )
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    add_title(slide, "3. 모델이 중요하게 본 신호", "부하, 마모, 회전 조건이 고장 리스크를 설명")
    slide.shapes.add_picture(str(FEATURE_IMG_PATH), Inches(0.75), Inches(1.2), width=Inches(6.0), height=Inches(4.1))
    top_features = feature.head(5).copy()
    top_features["feature"] = top_features["feature"].str.replace("num__", "", regex=False).str.replace("cat__", "", regex=False)
    for i, row in enumerate(top_features.itertuples(index=False), start=0):
        add_card(slide, 7.15 + (i % 2) * 2.45, 1.35 + (i // 2) * 1.35, 2.1, 1.0, short_feature_label(row.feature), str(int(row.importance)), "importance")
    add_text(slide, "Insight", 7.15, 5.5, 1.2, 0.3, size=13, color=LG_RED, bold=True)
    add_text(slide, "power_proxy와 공구마모가 상위 신호입니다. 단순 센서값보다 회전속도×토크 같은 조합 feature가 고장 위험을 더 잘 설명합니다.", 8.05, 5.45, 4.3, 0.75, size=12, color=LG_DARK)
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_bg(slide, LG_DARK)
    add_title(slide, "4. 그래도 모델이 놓치는 케이스", "False Negative와 margin이 큰 오분류가 다음 개선 포인트", dark=True)
    add_card(slide, 0.9, 1.35, 1.6, 1.05, "False Negative", str(fn), "missed failure", fill=WHITE)
    add_card(slide, 2.75, 1.35, 1.6, 1.05, "False Positive", str(fp), "false alarm", fill=WHITE)
    add_card(slide, 4.6, 1.35, 1.6, 1.05, "OOF rows", f"{len(best_oof):,}", "5-fold", fill=WHITE)
    add_text(slide, "아쉽게 틀린 케이스", 0.9, 2.75, 3.2, 0.3, size=15, color=PINK, bold=True)
    add_table(slide, near[["UDI", "actual", "predicted_probability", "Failure Type"]].round(3), 0.9, 3.15, 5.4, 1.25, font_size=7)
    add_text(slide, "고확신 오분류 케이스", 6.95, 2.75, 3.3, 0.3, size=15, color=PINK, bold=True)
    add_table(slide, high[["UDI", "actual", "predicted_probability", "Failure Type"]].round(3), 6.95, 3.15, 5.4, 1.25, font_size=7)
    add_text(slide, "Label risk", 0.9, 5.35, 1.3, 0.3, size=13, color=PINK, bold=True)
    add_text(slide, "Target과 Failure Type 간 불일치 27건은 라벨 정책 검토 대상입니다. 고장 미탐 비용이 크다면 threshold를 0.5보다 낮추는 비용 기반 최적화가 필요합니다.", 2.0, 5.28, 10.4, 0.7, size=12, color=WHITE)
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_bg(slide, WHITE)
    add_title(slide, "5. 모델 고도화 방향", "오분류 분석을 운영 가능한 예지보전 정책으로 연결")
    steps = [
        ("01", "Threshold 정책", "FN 비용을 크게 둔 cost-based threshold로 운영 리스크 최소화"),
        ("02", "라벨 정책", "Target/Failure Type 불일치 27건을 제외·재라벨링·유지로 비교"),
        ("03", "2단계 모델", "1단계 고장 여부, 2단계 고장 유형 분류로 목적 분리"),
        ("04", "반복 검증", "Repeated CV와 OOF calibration으로 안정성 검증"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        x = 0.9 + (i % 2) * 6.0
        y = 1.45 + (i // 2) * 2.05
        add_card(slide, x, y, 0.85, 0.85, "", num, "", fill=PINK)
        add_text(slide, title, x + 1.05, y + 0.05, 4.2, 0.35, size=16, color=LG_RED, bold=True)
        add_text(slide, desc, x + 1.05, y + 0.48, 4.5, 0.55, size=12, color=LG_DARK)
    add_text(slide, "Key takeaway", 0.95, 6.35, 1.7, 0.3, size=14, color=LG_RED, bold=True)
    add_text(slide, "최고 모델은 이미 높은 F1을 보였지만, 실제 현장 적용은 False Negative 관리와 라벨 품질 검토가 성능 이상의 핵심입니다.", 2.35, 6.3, 9.6, 0.42, size=13, color=LG_DARK)
    add_footer(slide, 6)

    prs.save(PPTX_PATH)
    print(PPTX_PATH)


if __name__ == "__main__":
    create_deck()
